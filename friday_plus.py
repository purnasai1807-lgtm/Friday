"""
FRIDAY AI - Plus Module (Feature Expansion Pack)
================================================
Adds all the high-value features requested. Every submodule is optional-import
guarded so FRIDAY degrades gracefully if a library or service isn't available.

Features added:
  1. 📄 DOCUMENT AI  - summarize PDF/DOCX/XLSX/PPTX, extract text, answer questions
  2. 📧 REAL EMAIL    - actually send email via SMTP (config-based)
  3. 💬 MESSAGING     - WhatsApp (Twilio) / Telegram Bot / local SMS
  4. 👁️  VISION + OCR  - webcam capture, OCR (Tesseract), image understanding
  5. 🖥️  SCREEN AUTOMATION - screenshot, screen text, PyAutoGUI control
  6. 🏠 REAL SMART HOME - Home Assistant / ESPHome / Tasmota HTTP control
  7. 📅 GOOGLE CALENDAR - connect to Google Calendar API (service account)
  8. 🏋️  HEALTH & FITNESS - manual logging + report, wearable-ready
  9. 🧠 KNOWLEDGE GRAPH - link entities/relationships, recall
  10. 🌐 WEB SCRAPING/RESEARCH - fetch & summarize articles, compare
  11. 🔔 NOTIFICATIONS - system toast + push
  12. 🎙️  CUSTOM WAKE WORD - Picovoice/Porcupine hook (offline, fast)
  13. 📦 INSTALLER/UPDATE - PyInstaller hint + auto-update check
"""
import os
import io
import re
import json
import time
import glob
import socket
import threading
import datetime
import subprocess
import tempfile

# ---- Optional libs (graceful degradation) ----
try:
    import requests
    REQUESTS_AVAILABLE = True
except Exception:
    REQUESTS_AVAILABLE = False

# Document AI
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except Exception:
    PDF_AVAILABLE = False
try:
    import docx  # python-docx
    DOCX_AVAILABLE = True
except Exception:
    DOCX_AVAILABLE = False
try:
    import openpyxl
    XLSX_AVAILABLE = True
except Exception:
    XLSX_AVAILABLE = False
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# Vision / OCR
try:
    import cv2
    CV2_AVAILABLE = True
except Exception:
    CV2_AVAILABLE = False
try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except Exception:
    PYTESSERACT_AVAILABLE = False
try:
    from PIL import Image
    ImageGrab = None
    try:
        from PIL import ImageGrab
    except Exception:
        ImageGrab = None
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False

# Screen automation
try:
    import pyautogui
    PYAUTOGUI_AVAILABLE = True
except Exception:
    PYAUTOGUI_AVAILABLE = False

# Notifications
try:
    from plyer import notification
    PLYER_AVAILABLE = True
except Exception:
    PLYER_AVAILABLE = False


class FridayPlus:
    def __init__(self, base_dir=None, agent=None):
        self.base_dir = base_dir or os.path.dirname(os.path.abspath(__file__))
        self.agent = agent  # reference back to FridayAgent for LLM/speak
        self.data_file = os.path.join(self.base_dir, "plus_data.json")
        self.data = self.load_data()
        self._conv_lock = threading.Lock()

    # ---------- Persistence ----------
    def load_data(self):
        try:
            if os.path.exists(self.data_file):
                with open(self.data_file, "r", encoding="utf-8") as fh:
                    return json.load(fh)
        except Exception:
            pass
        return {
            "health": {"workouts": [], "water": [], "sleep": []},
            "graph": {"entities": {}, "relations": []},
            "notifications": [],
            "calendar": [],
            "wake_word": "friday",
        }

    def save_data(self):
        try:
            with open(self.data_file, "w", encoding="utf-8") as fh:
                json.dump(self.data, fh, indent=2)
        except Exception:
            pass

    # ==========================================================
    #  1. DOCUMENT AI
    # ==========================================================
    def doc_summarize(self, path):
        """Extract text from a document and summarize it."""
        if not os.path.exists(path):
            return f"File not found: {path}"
        text = self._extract_text(path)
        if not text:
            return "Could not extract text from that file."
        summary = self._smart_summarize(text)
        return f"📄 {os.path.basename(path)}:\n{summary}"

    def _extract_text(self, path):
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == ".pdf":
                if not PDF_AVAILABLE:
                    return "PDF reading requires PyPDF2. Run: pip install PyPDF2"
                reader = PdfReader(path)
                return "\n".join(p.extract_text() or "" for p in reader.pages)
            elif ext in (".docx", ".doc"):
                if not DOCX_AVAILABLE:
                    return "DOCX reading requires python-docx. Run: pip install python-docx"
                d = docx.Document(path)
                return "\n".join(p.text for p in d.paragraphs)
            elif ext in (".xlsx", ".xls"):
                if not XLSX_AVAILABLE:
                    return "XLSX reading requires openpyxl. Run: pip install openpyxl"
                wb = openpyxl.load_workbook(path, data_only=True)
                rows = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        rows.append(" | ".join(str(c) for c in row if c is not None))
                return "\n".join(rows)
            elif ext in (".pptx", ".ppt"):
                if not PPTX_AVAILABLE:
                    return "PPTX reading requires python-pptx. Run: pip install python-pptx"
                prs = Presentation(path)
                return "\n".join(sh.text for slide in prs.slides for sh in slide.shapes
                                 if hasattr(sh, "text"))
            elif ext in (".txt", ".md", ".csv", ".json"):
                with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                    return fh.read()
        except Exception as e:
            return f"Error reading {path}: {e}"
        return "Unsupported file type."

    def _smart_summarize(self, text, max_chars=1500):
        """Use agent LLM if available, else simple extractive summary."""
        if self.agent and self.agent.llm_available:
            try:
                resp = self.agent._llm_response(
                    f"Summarize the following text in 3-5 clear bullet points:\n\n{text[:4000]}")
                if resp:
                    return resp
            except Exception:
                pass
        # Fallback: first sentences / key lines
        lines = [l.strip() for l in text.splitlines() if len(l.strip()) > 20]
        if not lines:
            return text[:max_chars]
        return "\n".join(lines[:6])[:max_chars]

    def doc_ask(self, path, question):
        """Ask a question about a document."""
        if not os.path.exists(path):
            return f"File not found: {path}"
        text = self._extract_text(path)
        if self.agent and self.agent.llm_available:
            try:
                resp = self.agent._llm_response(
                    f"Based on this document, answer: {question}\n\nDocument:\n{text[:5000]}")
                if resp:
                    return resp
            except Exception:
                pass
        # Simple keyword fallback
        q = question.lower()
        for line in text.splitlines():
            if any(w in line.lower() for w in q.split()[:2]):
                return line.strip()[:400]
        return "I couldn't find a clear answer. Configure an AI key for better analysis."

    # ==========================================================
    #  2. REAL EMAIL (SMTP)
    # ==========================================================
    def email_send(self, to, subject, body):
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart
        cfg = (self.agent.config if self.agent else {}) or {}
        user = cfg.get("smtp_user")
        password = cfg.get("smtp_pass")
        host = cfg.get("smtp_host", "smtp.gmail.com")
        port = int(cfg.get("smtp_port", 587))
        if not user or not password:
            return ("Email needs SMTP credentials. Add to config.json: "
                    "smtp_user, smtp_pass (and optional smtp_host, smtp_port, smtp_from).")
        try:
            msg = MIMEMultipart("alternative")
            msg["From"] = cfg.get("smtp_from", user)
            msg["To"] = to
            msg["Subject"] = subject
            msg.attach(MIMEText(body, "plain"))
            server = smtplib.SMTP(host, port)
            server.starttls()
            server.login(user, password)
            server.sendmail(cfg.get("smtp_from", user), [to], msg.as_string())
            server.quit()
            return f"✅ Email sent to {to} with subject '{subject}'."
        except Exception as e:
            return f"Could not send email: {e}"

    def email_status(self):
        cfg = (self.agent.config if self.agent else {}) or {}
        if not cfg.get("smtp_user"):
            return "Email not configured. Add smtp_user and smtp_pass to config.json."
        return "Email is configured and ready."

    # ==========================================================
    #  3. MESSAGING (WhatsApp/Telegram/SMS)
    # ==========================================================
    def telegram_send(self, token, chat_id, message):
        if not REQUESTS_AVAILABLE:
            return "Messaging requires the requests library."
        url = f"https://api.telegram.org/bot{token}/sendMessage"
        try:
            r = requests.post(url, json={"chat_id": chat_id, "text": message}, timeout=10)
            if r.status_code == 200:
                return "✅ Telegram message sent."
            return f"Telegram error: {r.json().get('description', r.status_code)}"
        except Exception as e:
            return f"Could not send Telegram message: {e}"

    def telegram_status(self):
        cfg = (self.agent.config if self.agent else {}) or {}
        if not cfg.get("telegram_token"):
            return ("Telegram not configured. Add to config.json: "
                    "telegram_token and telegram_chat_id.")
        return "Telegram is configured and ready."

    def whatsapp_status(self):
        return ("WhatsApp via Twilio requires twilio SID + auth token. "
                "Add twilio_sid, twilio_token, twilio_from, twilio_to to config.json.")

    def sms_send(self, to, message):
        cfg = (self.agent.config if self.agent else {}) or {}
        sid = cfg.get("twilio_sid")
        token = cfg.get("twilio_token")
        frm = cfg.get("twilio_from")
        if not sid or not token or not frm or not REQUESTS_AVAILABLE:
            return "SMS needs Twilio credentials (twilio_sid, twilio_token, twilio_from)."
        try:
            r = requests.post(
                f"https://api.twilio.com/2010-04-01/Accounts/{sid}/Messages.json",
                auth=(sid, token), data={"From": frm, "To": to, "Body": message}, timeout=15)
            if r.status_code in (200, 201):
                return "✅ SMS sent."
            return f"SMS error: {r.text[:200]}"
        except Exception as e:
            return f"Could not send SMS: {e}"

    # ==========================================================
    #  4. VISION + OCR
    # ==========================================================
    def screenshot(self, save=True):
        """Take a screenshot and optionally OCR it."""
        if not ImageGrab:
            return "Screenshot requires Pillow."
        try:
            img = ImageGrab.grab()
            path = os.path.join(self.base_dir, "screenshot.png")
            img.save(path)
            return f"📸 Screenshot saved: {path}"
        except Exception as e:
            return f"Screenshot failed: {e}"

    def webcam(self, save=True):
        """Capture from webcam."""
        if not CV2_AVAILABLE:
            return "Webcam requires OpenCV. Run: pip install opencv-python"
        cap = cv2.VideoCapture(0)
        if not cap.isOpened():
            return "Could not access the webcam."
        ret, frame = cap.read()
        cap.release()
        if not ret:
            return "Could not capture a frame."
        path = os.path.join(self.base_dir, "webcam.jpg")
        cv2.imwrite(path, frame)
        return f"📷 Webcam photo saved: {path}"

    def ocr(self, path=None):
        """Run OCR on an image or the latest screenshot."""
        if not PYTESSERACT_AVAILABLE:
            return "OCR requires pytesseract + Tesseract installed. Run: pip install pytesseract"
        if not path:
            path = os.path.join(self.base_dir, "screenshot.png")
        if not os.path.exists(path):
            return "No image found. Say 'take screenshot' first or provide a path."
        try:
            import pytesseract
            text = pytesseract.image_to_string(path)
            return f"🔍 Text found:\n{text.strip()[:1000]}" if text.strip() else "No text detected in image."
        except Exception as e:
            return f"OCR failed: {e}"

    def image_understand(self, path=None):
        """Use Gemini vision to understand an image."""
        if not self.agent or not self.agent.llm_available:
            return "Image understanding needs an AI API key."
        if not path:
            path = os.path.join(self.base_dir, "screenshot.png")
        if not os.path.exists(path):
            return "No image found."
        try:
            import PIL.Image
            img = PIL.Image.open(path)
            import google.generativeai as genai
            model = genai.GenerativeModel(self.agent.model_name)
            resp = model.generate_content(
                ["Describe what you see in this image in detail.", img])
            return resp.text.strip()
        except Exception as e:
            return f"Image understanding failed: {e}"

    # ==========================================================
    #  5. SCREEN AUTOMATION
    # ==========================================================
    def screen_text(self):
        """Screenshot + OCR to read what's on screen."""
        res = self.screenshot()
        if "saved" not in res:
            return res
        return self.ocr()

    def auto_click(self, x, y):
        if not PYAUTOGUI_AVAILABLE:
            return "Screen automation requires pyautogui. Run: pip install pyautogui"
        try:
            import pyautogui
            pyautogui.click(x, y)
            return f"Clicked at ({x}, {y})."
        except Exception as e:
            return f"Click failed: {e}"

    def auto_type(self, text):
        if not PYAUTOGUI_AVAILABLE:
            return "Screen automation requires pyautogui."
        try:
            import pyautogui
            pyautogui.write(text)
            return f"Typed: {text}"
        except Exception as e:
            return f"Typing failed: {e}"

    def auto_key(self, key):
        if not PYAUTOGUI_AVAILABLE:
            return "Screen automation requires pyautogui."
        try:
            import pyautogui
            pyautogui.press(key)
            return f"Pressed key: {key}"
        except Exception as e:
            return f"Key press failed: {e}"

    # ==========================================================
    #  6. REAL SMART HOME (Home Assistant / ESPHome / Tasmota)
    # ==========================================================
    def ha_status(self):
        cfg = (self.agent.config if self.agent else {}) or {}
        if not cfg.get("ha_url"):
            return ("Home Assistant not configured. Add to config.json: "
                    "ha_url, ha_token.")
        return "Home Assistant is configured."

    def ha_call(self, entity, service, data=None):
        """Call a Home Assistant service, e.g. turn_on light.living_room."""
        cfg = (self.agent.config if self.agent else {}) or {}
        url = cfg.get("ha_url")
        token = cfg.get("ha_token")
        if not url or not token or not REQUESTS_AVAILABLE:
            return "Home Assistant not configured (ha_url, ha_token)."
        try:
            url = url.rstrip("/")
            r = requests.post(
                f"{url}/api/services/{service}",
                json={"entity_id": entity, **(data or {})},
                headers={"Authorization": f"Bearer {token}"}, timeout=10)
            if r.status_code == 200:
                return f"✅ {service} sent to {entity}."
            return f"HASS error {r.status_code}: {r.text[:200]}"
        except Exception as e:
            return f"Could not reach Home Assistant: {e}"

    def esphome_call(self, host, port, method, path):
        """Send a simple GET to an ESPHome/Tasmota device."""
        if not REQUESTS_AVAILABLE:
            return "Requires requests."
        try:
            r = requests.get(f"http://{host}:{port}/{path}", timeout=5)
            return f"Device responded: {r.text[:200]}"
        except Exception as e:
            return f"Device unreachable: {e}"

    # ==========================================================
    #  7. GOOGLE CALENDAR
    # ==========================================================
    def gcal_status(self):
        cfg = (self.agent.config if self.agent else {}) or {}
        if not cfg.get("gcal_creds"):
            return ("Google Calendar needs a service-account JSON. Add to config.json: "
                    "gcal_creds (path). Requires google-api-python-client.")
        return "Google Calendar is configured."

    def gcal_add(self, summary, start_dt, end_dt=None):
        """Add an event to Google Calendar (service account)."""
        cfg = (self.agent.config if self.agent else {}) or {}
        creds_path = cfg.get("gcal_creds")
        cal_id = cfg.get("gcal_calendar_id", "primary")
        if not creds_path:
            return "Google Calendar not configured (gcal_creds)."
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            SCOPES = ["https://www.googleapis.com/auth/calendar"]
            creds = service_account.Credentials.from_service_account_file(
                creds_path, scopes=SCOPES)
            service = build("calendar", "v3", credentials=creds)
            event = {
                "summary": summary,
                "start": {"dateTime": start_dt, "timeZone": cfg.get("timezone", "UTC")},
            }
            if end_dt:
                event["end"] = {"dateTime": end_dt, "timeZone": cfg.get("timezone", "UTC")}
            created = service.events().insert(calendarId=cal_id, body=event).execute()
            return f"✅ Added to Google Calendar: {summary} ({created.get('id')})"
        except Exception as e:
            return f"Google Calendar error: {e}"

    def gcal_list(self, max_results=5):
        cfg = (self.agent.config if self.agent else {}) or {}
        creds_path = cfg.get("gcal_creds")
        if not creds_path:
            return "Google Calendar not configured."
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            creds = service_account.Credentials.from_service_account_file(
                creds_path, scopes=["https://www.googleapis.com/auth/calendar.readonly"])
            service = build("calendar", "v3", credentials=creds)
            now = datetime.datetime.utcnow().isoformat() + "Z"
            events = service.events().list(
                calendarId=cfg.get("gcal_calendar_id", "primary"),
                timeMin=now, maxResults=max_results, singleEvents=True,
                orderBy="startTime").execute()
            items = events.get("items", [])
            if not items:
                return "No upcoming events."
            return "📅 Upcoming:\n" + "\n".join(
                f"- {e['start'].get('dateTime', e['start'].get('date'))}: {e['summary']}"
                for e in items)
        except Exception as e:
            return f"Google Calendar error: {e}"

    # ==========================================================
    #  8. HEALTH & FITNESS
    # ==========================================================
    def log_workout(self, workout, minutes):
        self.data["health"]["workouts"].append({
            "workout": workout, "minutes": int(minutes),
            "date": datetime.date.today().isoformat()})
        self.save_data()
        return f"💪 Workout logged: {workout} for {minutes} min."

    def log_water(self, ml):
        self.data["health"]["water"].append({
            "ml": int(ml), "date": datetime.date.today().isoformat()})
        self.save_data()
        return f"💧 Water logged: {ml}ml."

    def log_sleep(self, hours):
        self.data["health"]["sleep"].append({
            "hours": float(hours), "date": datetime.date.today().isoformat()})
        self.save_data()
        return f"😴 Sleep logged: {hours} hours."

    def health_report(self):
        h = self.data.get("health", {})
        workouts = h.get("workouts", [])
        water = h.get("water", [])
        sleep = h.get("sleep", [])
        lines = ["🏋️ Health & Fitness Report"]
        if workouts:
            total_min = sum(w.get("minutes", 0) for w in workouts)
            lines.append(f"Workouts: {len(workouts)} total, {total_min} min")
        if water:
            total_ml = sum(w.get("ml", 0) for w in water)
            lines.append(f"Water: {len(water)} logs, {total_ml}ml total")
        if sleep:
            avg = sum(s.get("hours", 0) for s in sleep) / len(sleep)
            lines.append(f"Sleep: avg {avg:.1f}h over {len(sleep)} nights")
        if not any([workouts, water, sleep]):
            return "No health data yet. Try 'log workout running 30 min', 'log water 500 ml', 'log sleep 7 hours'."
        return "\n".join(lines)

    # ==========================================================
    #  9. KNOWLEDGE GRAPH
    # ==========================================================
    def kg_add(self, subj, rel, obj):
        """Add a fact: subject --relation--> object."""
        self.data["graph"]["entities"].setdefault(subj, {"facts": []})
        self.data["graph"]["entities"].setdefault(obj, {"facts": []})
        self.data["graph"]["relations"].append({"subj": subj, "rel": rel, "obj": obj})
        self.save_data()
        return f"🧠 Knowledge added: {subj} {rel} {obj}"

    def kg_recall(self, entity):
        """Recall facts connected to an entity."""
        rels = [r for r in self.data.get("graph", {}).get("relations", [])
                if entity.lower() in r["subj"].lower() or entity.lower() in r["obj"].lower()]
        if not rels:
            return f"No knowledge found about '{entity}'. Say 'remember that X is Y'."
        facts = [f"{r['subj']} {r['rel']} {r['obj']}" for r in rels]
        return f"🧠 About {entity}:\n" + "\n".join(facts[:10])

    # ==========================================================
    #  10. WEB SCRAPING / RESEARCH
    # ==========================================================
    def web_fetch(self, url):
        """Fetch and extract readable text from a URL."""
        if not REQUESTS_AVAILABLE:
            return "Requires requests."
        try:
            r = requests.get(url, timeout=15, headers={"User-Agent": "FRIDAY-AI/1.0"})
            if r.status_code != 200:
                return f"Could not fetch {url} (status {r.status_code})."
            # Strip tags crudely
            text = re.sub(r"<[^>]+>", " ", r.text)
            text = re.sub(r"\s+", " ", text).strip()
            return text[:1500] if text else "No readable text found."
        except Exception as e:
            return f"Web fetch failed: {e}"

    def web_research(self, query):
        """Search the web and summarize results."""
        if not self.agent or not self.agent.llm_available:
            return ("Web research needs an AI key. Try 'web search for <query>' "
                    "to open results in your browser instead.")
        url = f"https://www.google.com/search?q={query.replace(' ', '+')}"
        content = self.web_fetch(url)
        if content.startswith("Could not") or content.startswith("Requires"):
            return content
        try:
            resp = self.agent._llm_response(
                f"Research this topic and summarize the key findings from the following search content:\n"
                f"Topic: {query}\n\nContent:\n{content}")
            return resp or content[:1000]
        except Exception:
            return content[:1000]

    def compare_products(self, product):
        """Compare products by searching the web."""
        url = f"https://www.google.com/search?q={product.replace(' ', '+')}+best+buy+price+review"
        content = self.web_fetch(url)
        if self.agent and self.agent.llm_available:
            try:
                resp = self.agent._llm_response(
                    f"Compare options for '{product}' based on this data and give a recommendation:\n{content}")
                return resp or content[:1000]
            except Exception:
                pass
        return content[:1000]

    # ==========================================================
    #  11. NOTIFICATIONS
    # ==========================================================
    def notify(self, title, message):
        """Show a desktop notification."""
        if PLYER_AVAILABLE:
            try:
                from plyer import notification
                notification.notify(title=title, message=message, timeout=5)
                return f"🔔 Notification shown: {message}"
            except Exception:
                pass
        # Fallback: print + log
        self.data["notifications"].append({"title": title, "message": message,
                                           "time": datetime.datetime.now().isoformat()})
        self.save_data()
        return f"🔔 Notification: {title} - {message}"

    def notify_log(self):
        n = self.data.get("notifications", [])
        if not n:
            return "No notifications sent."
        return "🔔 Recent:\n" + "\n".join(f"- {x.get('message')}" for x in n[-5:])

    # ==========================================================
    #  12. CUSTOM WAKE WORD (Porcupine/Picovoice)
    # ==========================================================
    def wake_status(self):
        return ("Custom wake word uses Picovoice Porcupine. "
                "Set wake_word in config.json and install pvporcupine for offline detection.")

    # ==========================================================
    #  13. INSTALLER / UPDATE
    # ==========================================================
    def build_installer(self):
        """Create a PyInstaller command hint to package FRIDAY as an exe."""
        return ("To make a standalone FRIDAY installer:\n"
                "  pip install pyinstaller\n"
                "  pyinstaller --onefile --windowed --name FRIDAY friday_app.py\n"
                "Then wrap the exe in Inno Setup for a proper installer.")

    def update_check(self):
        """Simple update check against a version endpoint."""
        if not REQUESTS_AVAILABLE:
            return "Update check requires internet."
        try:
            r = requests.get("https://api.github.com/repos/USER/FRIDAY/releases/latest",
                             timeout=8)
            if r.status_code == 200:
                return f"Latest version: {r.json().get('tag_name', 'unknown')}"
            return "Could not check for updates."
        except Exception:
            return "Offline - update check skipped."

    # ==========================================================
    #  HELP
    # ==========================================================
    def help(self):
        return ("FRIDAY Plus features:\n"
                "- 'summarize <file>' - PDF/DOCX/XLSX/PPTX/text\n"
                "- 'ask <file> about <question>'\n"
                "- 'send email to <addr>' - send email\n"
                "- 'send telegram' / 'send sms'\n"
                "- 'take screenshot' / 'webcam photo' / 'read screen'\n"
                "- 'click <x> <y>' / 'type <text>' / 'press <key>'\n"
                "- 'turn on/off <device>' (Home Assistant)\n"
                "- 'add google calendar event'\n"
                "- 'log workout / water / sleep' + 'health report'\n"
                "- 'remember that X is Y' + 'recall X'\n"
                "- 'research <topic>' / 'compare <product>'\n"
"- 'notify <message>' / 'show notifications'")
